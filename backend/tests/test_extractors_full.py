"""
Full edge-case test suite for all extractors.
Run: python tests/test_extractors_full.py
"""
import sys
import tempfile
import traceback
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))

from app.shared.utils.extractors.pdf_extractor  import extract_pdf
from app.shared.utils.extractors.docx_extractor import extract_docx
from app.shared.utils.extractors.xlsx_extractor import extract_xlsx
from app.shared.utils.extractors.pptx_extractor import extract_pptx

PASS = "[PASS]"
FAIL = "[FAIL]"

results = []

def test(name, fn):
    try:
        fn()
        print(f"  {PASS} {name}")
        results.append((name, True, None))
    except Exception as e:
        print(f"  {FAIL} {name}")
        print(f"       {type(e).__name__}: {e}")
        results.append((name, False, e))


def section(title):
    print(f"\n{'='*55}")
    print(f"  {title}")
    print(f"{'='*55}")


# â”€â”€â”€ DOCX fixtures â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

def make_docx_empty(path):
    from docx import Document
    Document().save(path)

def make_docx_basic(path):
    from docx import Document
    doc = Document()
    doc.add_paragraph("Hello RAG")
    doc.add_paragraph("Second line")
    t = doc.add_table(rows=2, cols=2)
    t.cell(0,0).text = "Name"; t.cell(0,1).text = "Score"
    t.cell(1,0).text = "An";   t.cell(1,1).text = "100"
    doc.save(path)

def make_docx_merged_cells(path):
    from docx import Document
    from docx.oxml.ns import qn
    import copy
    doc = Document()
    t = doc.add_table(rows=2, cols=3)
    t.cell(0,0).text = "A"; t.cell(0,1).text = "B"; t.cell(0,2).text = "C"
    # merge col 0 and col 1 in row 1
    t.cell(1,0).merge(t.cell(1,1))
    t.cell(1,0).text = "merged"
    t.cell(1,2).text = "solo"
    doc.save(path)

def make_docx_with_image(path, image_path):
    from docx import Document
    from docx.shared import Inches
    doc = Document()
    doc.add_paragraph("Doc with image")
    doc.add_picture(str(image_path), width=Inches(1))
    doc.save(path)

def make_tiny_png(path):
    import struct, zlib
    def chunk(tag, data):
        c = struct.pack(">I", len(data)) + tag + data
        return c + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF)
    sig = b"\x89PNG\r\n\x1a\n"
    ihdr = chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0))
    raw  = b"\x00\xff\x00\x00"
    idat = chunk(b"IDAT", zlib.compress(raw))
    iend = chunk(b"IEND", b"")
    path.write_bytes(sig + ihdr + idat + iend)

# â”€â”€â”€ XLSX fixtures â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

def make_xlsx_basic(path):
    from openpyxl import Workbook
    wb = Workbook(); ws = wb.active; ws.title = "Data"
    ws.append(["Name", "Score"]); ws.append(["An", 100]); ws.append(["Binh", 85])
    wb.save(path)

def make_xlsx_empty_sheet(path):
    from openpyxl import Workbook
    wb = Workbook(); ws = wb.active; ws.title = "Empty"
    wb.save(path)

def make_xlsx_multi_sheet(path):
    from openpyxl import Workbook
    wb = Workbook()
    ws1 = wb.active; ws1.title = "Sheet1"
    ws1.append(["X", "Y"]); ws1.append([1, 2])
    ws2 = wb.create_sheet("Sheet2")
    ws2.append(["A", "B"]); ws2.append([3, 4])
    wb.save(path)

def make_xlsx_none_values(path):
    from openpyxl import Workbook
    wb = Workbook(); ws = wb.active; ws.title = "Nulls"
    ws.append(["Name", "Value", "Note"])
    ws.append(["An", None, "ok"])
    ws.append([None, 99, None])
    wb.save(path)

def make_xlsx_duplicate_headers(path):
    from openpyxl import Workbook
    wb = Workbook(); ws = wb.active; ws.title = "Dup"
    ws.append(["Name", "Name", "Score"])
    ws.append(["An", "Nguyen", 100])
    wb.save(path)

def make_xlsx_none_header(path):
    from openpyxl import Workbook
    wb = Workbook(); ws = wb.active; ws.title = "NoneHdr"
    ws.append([None, "Score"])
    ws.append(["An", 100])
    wb.save(path)

# â”€â”€â”€ PPTX fixtures â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

def make_pptx_basic(path):
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Title"
    slide.placeholders[1].text = "Body text"
    prs.save(path)

def make_pptx_empty_slide(path):
    from pptx import Presentation
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    prs.save(path)

def make_pptx_multi_slide(path):
    from pptx import Presentation
    prs = Presentation()
    for i in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Slide {i+1}"
        slide.placeholders[1].text = f"Content {i+1}"
    prs.save(path)

def make_pptx_with_table(path):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rows, cols = 3, 2
    tbl = slide.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(4), Inches(2)).table
    tbl.cell(0,0).text = "Name"; tbl.cell(0,1).text = "Score"
    tbl.cell(1,0).text = "An";   tbl.cell(1,1).text = "100"
    tbl.cell(2,0).text = "Binh"; tbl.cell(2,1).text = "85"
    prs.save(path)

# â”€â”€â”€ PDF fixtures â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

def make_pdf_basic(path):
    from fpdf import FPDF
    pdf = FPDF(); pdf.add_page(); pdf.set_font("Helvetica", size=12)
    pdf.cell(0, 10, "Hello PDF", new_x="LMARGIN", new_y="NEXT")
    pdf.output(str(path))

def make_pdf_multipage(path):
    from fpdf import FPDF
    pdf = FPDF(); pdf.set_font("Helvetica", size=12)
    for i in range(3):
        pdf.add_page()
        pdf.cell(0, 10, f"Page {i+1} content", new_x="LMARGIN", new_y="NEXT")
    pdf.output(str(path))

def make_pdf_empty(path):
    from fpdf import FPDF
    pdf = FPDF(); pdf.add_page()
    pdf.output(str(path))


# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# DOCX tests
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•

section("DOCX extractor")

with tempfile.TemporaryDirectory() as tmp:
    tmp = Path(tmp)
    img_dir = tmp / "img"; img_dir.mkdir()

    def t_docx_empty():
        p = tmp / "empty.docx"; make_docx_empty(p)
        r = extract_docx(str(p), str(img_dir))
        assert r.text == "", f"expected empty text, got {repr(r.text)}"
        assert r.tables == []
        assert r.images == []

    def t_docx_text_and_table():
        p = tmp / "basic.docx"; make_docx_basic(p)
        r = extract_docx(str(p), str(img_dir))
        assert "Hello RAG" in r.text
        assert len(r.tables) == 1
        assert r.tables[0]["data"][0]["Name"] == "An"
        assert r.tables[0]["data"][0]["Score"] == "100"

    def t_docx_merged_cells_no_duplicate():
        p = tmp / "merged.docx"; make_docx_merged_cells(p)
        r = extract_docx(str(p), str(img_dir))
        assert len(r.tables) == 1
        row = r.tables[0]["data"][0]
        # merged cell should appear once, not duplicated as a key
        assert list(row.values()).count("merged") == 1, f"duplicate merged cell: {row}"

    def t_docx_image_extracted():
        png = tmp / "tiny.png"; make_tiny_png(png)
        p   = tmp / "with_img.docx"; make_docx_with_image(p, png)
        r   = extract_docx(str(p), str(img_dir))
        assert len(r.images) >= 1
        assert all(Path(ip).exists() for ip in r.images)

    def t_docx_metadata_has_source():
        p = tmp / "meta.docx"; make_docx_basic(p)
        r = extract_docx(str(p), str(img_dir))
        assert "source" in r.metadata

    test("empty doc - empty result",        t_docx_empty)
    test("text + table extracted",          t_docx_text_and_table)
    test("merged cells not duplicated",     t_docx_merged_cells_no_duplicate)
    test("embedded image saved to disk",    t_docx_image_extracted)
    test("metadata contains source",        t_docx_metadata_has_source)


# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# XLSX tests
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•

section("XLSX extractor")

with tempfile.TemporaryDirectory() as tmp:
    tmp = Path(tmp)

    def t_xlsx_basic():
        p = tmp / "basic.xlsx"; make_xlsx_basic(p)
        r = extract_xlsx(str(p))
        assert "An" in r.text
        assert len(r.tables) == 1
        assert r.tables[0]["data"][0]["Name"] == "An"
        assert r.tables[0]["data"][0]["Score"] == 100

    def t_xlsx_empty_sheet_skipped():
        p = tmp / "empty.xlsx"; make_xlsx_empty_sheet(p)
        r = extract_xlsx(str(p))
        assert r.tables == []
        assert r.text == ""

    def t_xlsx_multi_sheet():
        p = tmp / "multi.xlsx"; make_xlsx_multi_sheet(p)
        r = extract_xlsx(str(p))
        assert len(r.tables) == 2
        sheet_names = {t["sheet"] for t in r.tables}
        assert "Sheet1" in sheet_names
        assert "Sheet2" in sheet_names

    def t_xlsx_none_values_handled():
        p = tmp / "nulls.xlsx"; make_xlsx_none_values(p)
        r = extract_xlsx(str(p))
        assert len(r.tables) == 1
        rows = r.tables[0]["data"]
        # None values should be preserved (not crash)
        assert any(row.get("Value") is None for row in rows)

    def t_xlsx_none_header_fallback():
        p = tmp / "nonehdr.xlsx"; make_xlsx_none_header(p)
        r = extract_xlsx(str(p))
        assert len(r.tables) == 1
        # header None â†’ should fall back to col_0, not key "None"
        row = r.tables[0]["data"][0]
        assert "None" not in row, f"None used as key: {row}"

    def t_xlsx_duplicate_headers():
        p = tmp / "dup.xlsx"; make_xlsx_duplicate_headers(p)
        r = extract_xlsx(str(p))
        # duplicate header â†’ second key overwrites first â†’ "Name" == "Nguyen"
        # should either dedup or document the behavior â€” test it doesn't crash
        assert len(r.tables) == 1

    def t_xlsx_metadata():
        p = tmp / "meta.xlsx"; make_xlsx_basic(p)
        r = extract_xlsx(str(p))
        assert "sheets" in r.metadata
        assert "source" in r.metadata

    test("basic text + table",              t_xlsx_basic)
    test("empty sheet skipped",             t_xlsx_empty_sheet_skipped)
    test("multi-sheet â†’ multiple tables",   t_xlsx_multi_sheet)
    test("None values in cells OK",         t_xlsx_none_values_handled)
    test("None header â†’ col_N fallback",    t_xlsx_none_header_fallback)
    test("duplicate headers no crash",      t_xlsx_duplicate_headers)
    test("metadata contains sheets+source", t_xlsx_metadata)


# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# PPTX tests
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•

section("PPTX extractor")

with tempfile.TemporaryDirectory() as tmp:
    tmp = Path(tmp)
    img_dir = tmp / "img"; img_dir.mkdir()

    def t_pptx_basic():
        p = tmp / "basic.pptx"; make_pptx_basic(p)
        r = extract_pptx(str(p), str(img_dir))
        assert "Title" in r.text
        assert "Body text" in r.text

    def t_pptx_empty_slide_no_crash():
        p = tmp / "empty.pptx"; make_pptx_empty_slide(p)
        r = extract_pptx(str(p), str(img_dir))
        assert r.text == ""
        assert r.tables == []

    def t_pptx_multi_slide_order():
        p = tmp / "multi.pptx"; make_pptx_multi_slide(p)
        r = extract_pptx(str(p), str(img_dir))
        assert "[Slide 1]" in r.text
        assert "[Slide 2]" in r.text
        assert "[Slide 3]" in r.text
        assert r.metadata["slides"] == 3

    def t_pptx_table_extracted():
        p = tmp / "table.pptx"; make_pptx_with_table(p)
        r = extract_pptx(str(p), str(img_dir))
        assert len(r.tables) == 1
        assert r.tables[0]["data"][0]["Name"] == "An"
        assert r.tables[0]["data"][0]["Score"] == "100"

    def t_pptx_metadata():
        p = tmp / "meta.pptx"; make_pptx_basic(p)
        r = extract_pptx(str(p), str(img_dir))
        assert "slides" in r.metadata
        assert "source" in r.metadata

    test("basic text extracted",            t_pptx_basic)
    test("empty slide no crash",            t_pptx_empty_slide_no_crash)
    test("multi-slide order preserved",     t_pptx_multi_slide_order)
    test("table in slide extracted",        t_pptx_table_extracted)
    test("metadata contains slides+source", t_pptx_metadata)


# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# PDF tests
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•

section("PDF extractor")

with tempfile.TemporaryDirectory() as tmp:
    tmp = Path(tmp)
    img_dir = tmp / "img"; img_dir.mkdir()

    def t_pdf_basic():
        p = tmp / "basic.pdf"; make_pdf_basic(p)
        r = extract_pdf(str(p), str(img_dir))
        assert "Hello PDF" in r.text

    def t_pdf_multipage_text_joined():
        p = tmp / "multi.pdf"; make_pdf_multipage(p)
        r = extract_pdf(str(p), str(img_dir))
        assert "Page 1" in r.text
        assert "Page 2" in r.text
        assert "Page 3" in r.text
        assert r.metadata["pages"] == 3

    def t_pdf_empty_no_crash():
        p = tmp / "empty.pdf"; make_pdf_empty(p)
        r = extract_pdf(str(p), str(img_dir))
        assert r.text == ""

    def t_pdf_no_resource_leak():
        import gc
        import fitz
        p = tmp / "leak.pdf"; make_pdf_basic(p)
        before = len([o for o in gc.get_objects() if isinstance(o, fitz.Document)])
        extract_pdf(str(p), str(img_dir))
        gc.collect()
        after = len([o for o in gc.get_objects() if isinstance(o, fitz.Document)])
        assert after <= before, f"fitz.Document leak: {before} â†’ {after}"

    def t_pdf_metadata_has_pages():
        p = tmp / "meta.pdf"; make_pdf_multipage(p)
        r = extract_pdf(str(p), str(img_dir))
        assert r.metadata["pages"] == 3
        assert "source" in r.metadata

    test("basic text extracted",            t_pdf_basic)
    test("multi-page text joined",          t_pdf_multipage_text_joined)
    test("empty PDF no crash",              t_pdf_empty_no_crash)
    test("no fitz.Document resource leak",  t_pdf_no_resource_leak)
    test("metadata has pages + source",     t_pdf_metadata_has_pages)


# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# Summary
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•

passed = sum(1 for _, ok, _ in results if ok)
failed = [(n, e) for n, ok, e in results if not ok]
total  = len(results)

print(f"\n{'='*55}")
print(f"  Result: {passed}/{total} passed")
if failed:
    print(f"\n  Failed cases:")
    for name, err in failed:
        print(f"    - {name}: {type(err).__name__}: {err}")
print(f"{'='*55}\n")

sys.exit(0 if not failed else 1)

