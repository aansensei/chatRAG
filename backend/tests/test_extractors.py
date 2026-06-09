"""
Quick test for all extractors. Generates fixture files programmatically, no real files needed.
Run: python tests/test_extractors.py
"""
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))


def make_docx(path: Path) -> None:
    from docx import Document
    from docx.oxml.ns import qn
    import docx.oxml as oxml

    doc = Document()
    doc.add_paragraph("Hello from DOCX")
    doc.add_paragraph("Second paragraph about RAG pipeline")

    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Score"
    table.cell(1, 0).text = "An"
    table.cell(1, 1).text = "100"

    doc.save(path)


def make_xlsx(path: Path) -> None:
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Name", "Score", "Grade"])
    ws.append(["An", 100, "A"])
    ws.append(["Binh", 85, "B"])
    wb.save(path)


def make_pptx(path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "RAG Pipeline"
    slide.placeholders[1].text = "Extract - Chunk - Embed - Store"
    prs.save(path)


def make_pdf(path: Path) -> None:
    from fpdf import FPDF
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=12)
    pdf.cell(0, 10, "Hello from PDF", ln=True)
    pdf.cell(0, 10, "RAG pipeline document", ln=True)
    pdf.output(str(path))


def run_test(label: str, extract_fn, file_path: Path, image_dir: Path):
    print(f"\n{'='*50}")
    print(f"  {label}")
    print(f"{'='*50}")
    try:
        result = extract_fn(str(file_path), str(image_dir))
        print(f"  text    : {repr(result.text[:80])}{'...' if len(result.text) > 80 else ''}")
        print(f"  tables  : {len(result.tables)} table(s)")
        for i, t in enumerate(result.tables):
            print(f"    table[{i}]: {t}")
        print(f"  images  : {len(result.images)} image(s)")
        print(f"  metadata: {result.metadata}")
        print(f"  [PASS]")
        return True
    except Exception as e:
        print(f"  [FAIL] {type(e).__name__}: {e}")
        return False


def run_xlsx_test(file_path: Path):
    print(f"\n{'='*50}")
    print(f"  XLSX extractor")
    print(f"{'='*50}")
    from app.shared.utils.extractors.xlsx_extractor import extract_xlsx
    try:
        result = extract_xlsx(str(file_path))
        print(f"  text    : {repr(result.text[:80])}{'...' if len(result.text) > 80 else ''}")
        print(f"  tables  : {len(result.tables)} table(s)")
        for i, t in enumerate(result.tables):
            print(f"    table[{i}]: {t}")
        print(f"  metadata: {result.metadata}")
        print(f"  [PASS]")
        return True
    except Exception as e:
        print(f"  [FAIL] {type(e).__name__}: {e}")
        return False


if __name__ == "__main__":
    from app.shared.utils.extractors.pdf_extractor import extract_pdf
    from app.shared.utils.extractors.docx_extractor import extract_docx
    from app.shared.utils.extractors.pptx_extractor import extract_pptx

    with tempfile.TemporaryDirectory() as tmp:
        tmp = Path(tmp)
        image_dir = tmp / "images"
        image_dir.mkdir()

        docx_path = tmp / "test.docx"
        xlsx_path = tmp / "test.xlsx"
        pptx_path = tmp / "test.pptx"
        pdf_path  = tmp / "test.pdf"

        make_docx(docx_path)
        make_xlsx(xlsx_path)
        make_pptx(pptx_path)
        make_pdf(pdf_path)

        results = []
        results.append(run_test("DOCX extractor", extract_docx, docx_path, image_dir))
        results.append(run_xlsx_test(xlsx_path))
        results.append(run_test("PPTX extractor", extract_pptx, pptx_path, image_dir))
        results.append(run_test("PDF extractor",  extract_pdf,  pdf_path,  image_dir))

        passed = sum(results)
        total  = len(results)
        print(f"\n{'='*50}")
        print(f"  Result: {passed}/{total} passed")
        print(f"{'='*50}\n")
        sys.exit(0 if passed == total else 1)
