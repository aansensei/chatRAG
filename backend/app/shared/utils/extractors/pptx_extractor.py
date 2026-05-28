from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from .base import ExtractResult


def extract_pptx(file_path: str, image_output_dir: str = "extracted_images") -> ExtractResult:
    path = Path(file_path)
    image_dir = Path(image_output_dir)
    image_dir.mkdir(parents=True, exist_ok=True)

    result = ExtractResult()
    prs = Presentation(file_path)
    text_parts = []
    tables = []
    image_paths = []

    for slide_index, slide in enumerate(prs.slides):
        slide_texts = []

        for shape in slide.shapes:
            # text
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        slide_texts.append(line)

            # table
            if shape.has_table:
                rows = [
                    [cell.text.strip() for cell in row.cells]
                    for row in shape.table.rows
                ]
                if rows:
                    header = rows[0]
                    data = [
                        {header[i]: cell for i, cell in enumerate(row)}
                        for row in rows[1:]
                    ]
                    tables.append({
                        "slide": slide_index + 1,
                        "data": data
                    })

            # image
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image = shape.image
                ext = image.ext
                save_path = image_dir / f"{path.stem}_slide{slide_index + 1}_{shape.shape_id}.{ext}"
                save_path.write_bytes(image.blob)
                image_paths.append(str(save_path))

        if slide_texts:
            text_parts.append(f"[Slide {slide_index + 1}]\n" + "\n".join(slide_texts))

    result.text = "\n\n".join(text_parts)
    result.tables = tables
    result.images = image_paths
    result.metadata = {"slides": len(prs.slides), "source": str(path)}

    return result
